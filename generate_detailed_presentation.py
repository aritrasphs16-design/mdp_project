"""
ChemoFilter Detailed Presentation Generator
Generates ChemoFilter_Detailed_Presentation.pptx with comprehensive
chemical term definitions for the Chemistry Department panel.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG      = RGBColor(10,  20,  45)   # midnight navy
TITLE_CYAN   = RGBColor(0,  188, 212)   # electric teal
GOLD         = RGBColor(212, 175,  55)  # amber gold
WHITE        = RGBColor(255, 255, 255)
LIGHT_GRAY   = RGBColor(200, 210, 230)
ACCENT_TEAL  = RGBColor(0,  229, 255)
DIM_BLUE     = RGBColor(30,  50,  90)   # slightly lighter navy for boxes
GREEN        = RGBColor( 72, 199, 142)
RED_WARN     = RGBColor(255, 107, 107)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = 6   # fully blank layout in most pptx themes


# ── Low-level XML helpers ────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height,
                 text="", font_size=24, bold=False, italic=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    """Add a styled textbox to a slide (single-paragraph helper)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return txBox


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a filled rectangle (optional border)."""
    shape = slide.shapes.add_shape(
        1,  # 1 = rectangle (MSO_SHAPE_TYPE.RECTANGLE)
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _multiline_textbox(slide, left, top, width, height,
                       paragraphs, base_size=22, wrap=True):
    """
    Add a textbox with multiple styled paragraphs.
    paragraphs = list of dicts:
        {text, size, bold, italic, color, align, space_before}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    first = True
    for para in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "space_before" in para:
            p.space_before = Pt(para["space_before"])
        run = p.add_run()
        run.text = para.get("text", "")
        run.font.size  = Pt(para.get("size", base_size))
        run.font.bold  = para.get("bold", False)
        run.font.italic = para.get("italic", False)
        run.font.color.rgb = para.get("color", WHITE)
        run.font.name  = "Arial"
    return txBox


# ── Reusable slide builders ──────────────────────────────────────────────────

def make_blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    _set_bg(slide, DARK_BG)
    return slide


def add_slide_header(slide, title, subtitle=None):
    """Teal title bar at the top of a slide."""
    _add_rect(slide,
              Inches(0), Inches(0), SLIDE_W, Inches(1.2),
              DIM_BLUE)
    _add_textbox(slide,
                 Inches(0.3), Inches(0.1),
                 Inches(12.5), Inches(0.8),
                 text=title, font_size=32, bold=True,
                 color=TITLE_CYAN, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_textbox(slide,
                     Inches(0.3), Inches(0.85),
                     Inches(12.5), Inches(0.35),
                     text=subtitle, font_size=17, italic=True,
                     color=GOLD, align=PP_ALIGN.LEFT)


def add_gold_divider(slide, y_inches):
    """Thin gold horizontal rule."""
    from pptx.util import Pt as _Pt
    line = slide.shapes.add_shape(1,
        Inches(0.3), Inches(y_inches),
        Inches(12.7), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()


def title_slide(title, subtitle, footer=""):
    """Full-bleed title slide."""
    slide = make_blank_slide()
    # Background gradient-ish with two rects
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK_BG)
    _add_rect(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(1.7),
              DIM_BLUE)
    # Gold accent bar
    _add_rect(slide, Inches(0), Inches(2.6), Inches(0.18), Inches(1.5),
              GOLD)
    # Title
    _multiline_textbox(slide,
        Inches(0.4), Inches(1.8),
        Inches(12.5), Inches(2),
        [
            {"text": title,    "size": 38, "bold": True,
             "color": TITLE_CYAN, "align": PP_ALIGN.LEFT},
        ])
    # Subtitle
    _multiline_textbox(slide,
        Inches(0.4), Inches(3.2),
        Inches(11), Inches(1.5),
        [{"text": subtitle, "size": 22, "color": LIGHT_GRAY,
          "align": PP_ALIGN.LEFT}])
    if footer:
        _add_textbox(slide,
                     Inches(0.3), Inches(6.7), Inches(12.5), Inches(0.5),
                     text=footer, font_size=14, color=GOLD,
                     align=PP_ALIGN.LEFT)
    return slide


def section_divider(section_num, section_title, description=""):
    """Interstitial section slide."""
    slide = make_blank_slide()
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK_BG)
    _add_rect(slide, Inches(0), Inches(3.0), SLIDE_W, Inches(1.8),
              DIM_BLUE)
    _add_rect(slide, Inches(0), Inches(3.0), Inches(0.25), Inches(1.8),
              TITLE_CYAN)
    _multiline_textbox(slide,
        Inches(0.5), Inches(2.6),
        Inches(12), Inches(0.6),
        [{"text": f"Slide Set {section_num}", "size": 18,
          "color": GOLD, "bold": True, "align": PP_ALIGN.LEFT}])
    _multiline_textbox(slide,
        Inches(0.4), Inches(3.05),
        Inches(12.3), Inches(1.7),
        [{"text": section_title, "size": 34, "bold": True,
          "color": TITLE_CYAN, "align": PP_ALIGN.LEFT},
         {"text": description,   "size": 20, "color": LIGHT_GRAY,
          "align": PP_ALIGN.LEFT, "space_before": 6}])
    return slide


def definition_slide(title, subtitle, terms):
    """
    Slide with a title bar and a list of term+definition boxes.
    terms = list of (term, definition) tuples.
    """
    slide = make_blank_slide()
    add_slide_header(slide, title, subtitle)
    add_gold_divider(slide, 1.22)

    y = 1.4
    row_h = (7.5 - 1.5) / max(len(terms), 1)
    row_h = min(row_h, 1.25)

    for term, defn in terms:
        # Term label box
        _add_rect(slide,
                  Inches(0.25), Inches(y),
                  Inches(3.0), Inches(row_h - 0.05),
                  DIM_BLUE, GOLD)
        _multiline_textbox(slide,
            Inches(0.3), Inches(y + 0.05),
            Inches(2.9), Inches(row_h - 0.1),
            [{"text": term, "size": 19, "bold": True,
              "color": GOLD, "align": PP_ALIGN.LEFT}])
        # Definition box
        _multiline_textbox(slide,
            Inches(3.4), Inches(y + 0.04),
            Inches(9.6), Inches(row_h - 0.1),
            [{"text": defn, "size": 17, "color": WHITE,
              "align": PP_ALIGN.LEFT}])
        y += row_h

    return slide


def bullet_slide(title, subtitle, bullets, bullet_color=WHITE,
                 bullet_size=20):
    """Standard bullet-point slide."""
    slide = make_blank_slide()
    add_slide_header(slide, title, subtitle)
    add_gold_divider(slide, 1.22)

    y = 1.4
    for bullet in bullets:
        _add_rect(slide,
                  Inches(0.3), Inches(y + 0.12),
                  Inches(0.15), Inches(0.15),
                  GOLD)
        _multiline_textbox(slide,
            Inches(0.55), Inches(y),
            Inches(12.4), Inches(0.65),
            [{"text": bullet, "size": bullet_size,
              "color": bullet_color, "align": PP_ALIGN.LEFT}])
        y += 0.72

    return slide


def two_column_slide(title, subtitle, left_items, right_items,
                     left_header="", right_header=""):
    """Two-column layout slide."""
    slide = make_blank_slide()
    add_slide_header(slide, title, subtitle)
    add_gold_divider(slide, 1.22)

    # Left column header
    if left_header:
        _add_rect(slide, Inches(0.25), Inches(1.32),
                  Inches(6.1), Inches(0.4), TITLE_CYAN)
        _add_textbox(slide, Inches(0.3), Inches(1.33),
                     Inches(6.0), Inches(0.38),
                     text=left_header, font_size=16, bold=True,
                     color=DARK_BG, align=PP_ALIGN.CENTER)
    if right_header:
        _add_rect(slide, Inches(6.85), Inches(1.32),
                  Inches(6.1), Inches(0.4), GOLD)
        _add_textbox(slide, Inches(6.9), Inches(1.33),
                     Inches(6.0), Inches(0.38),
                     text=right_header, font_size=16, bold=True,
                     color=DARK_BG, align=PP_ALIGN.CENTER)

    y_start = 1.82
    row_h = 0.68
    max_rows = max(len(left_items), len(right_items))
    for i in range(max_rows):
        y = y_start + i * row_h
        if i < len(left_items):
            _add_rect(slide, Inches(0.25), Inches(y),
                      Inches(6.1), Inches(row_h - 0.06),
                      DIM_BLUE)
            _multiline_textbox(slide,
                Inches(0.35), Inches(y + 0.04),
                Inches(5.9), Inches(row_h - 0.12),
                [{"text": left_items[i], "size": 17,
                  "color": WHITE, "align": PP_ALIGN.LEFT}])
        if i < len(right_items):
            _add_rect(slide, Inches(6.85), Inches(y),
                      Inches(6.1), Inches(row_h - 0.06),
                      DIM_BLUE)
            _multiline_textbox(slide,
                Inches(6.95), Inches(y + 0.04),
                Inches(5.9), Inches(row_h - 0.12),
                [{"text": right_items[i], "size": 17,
                  "color": WHITE, "align": PP_ALIGN.LEFT}])
    return slide


def thank_you_slide():
    slide = make_blank_slide()
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK_BG)
    _add_rect(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(2.0),
              DIM_BLUE)
    _multiline_textbox(slide,
        Inches(0.5), Inches(1.5), Inches(12), Inches(1.5),
        [{"text": "Thank You", "size": 52, "bold": True,
          "color": TITLE_CYAN, "align": PP_ALIGN.CENTER}])
    _multiline_textbox(slide,
        Inches(0.5), Inches(3.1), Inches(12), Inches(0.8),
        [{"text": "Questions & Discussion", "size": 28,
          "color": GOLD, "align": PP_ALIGN.CENTER}])
    _multiline_textbox(slide,
        Inches(0.5), Inches(4.0), Inches(12), Inches(1.2),
        [{"text": "ChemoFilter | Chemistry Department Panel Presentation",
          "size": 18, "color": LIGHT_GRAY, "align": PP_ALIGN.CENTER}])
    return slide


# ════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ════════════════════════════════════════════════════════════════════════════

# ── SLIDE 1 · Title ──────────────────────────────────────────────────────────
title_slide(
    "ChemoFilter",
    "An Advanced Computational Drug Discovery & Screening Platform\n"
    "Comprehensive Overview for the Chemistry Department Panel",
    "Powered by RDKit · Anthropic Claude AI · Streamlit"
)

# ── SLIDE 2 · Agenda ─────────────────────────────────────────────────────────
bullet_slide(
    "Presentation Agenda",
    "What we will cover today",
    [
        "1 · Introduction & Problem Statement",
        "2 · Platform Overview — What is ChemoFilter?",
        "3 · ADMET & Physicochemical Properties (with definitions)",
        "4 · Drug-likeness Filters — Lipinski, Muegge, PAINS, Brenk",
        "5 · Toxicity Indicators — hepatotoxicity, hERG, Ames, CYP450 …",
        "6 · Structural Analysis — SMILES, InChI, Scaffolds, Alerts",
        "7 · Core Capabilities  ·  8 · ChemoScore  ·  9 · Workflow",
        "10 · AI Integration  ·  11 · Unique Features  ·  12 · Future Scope",
        "13 · References  ·  14 · Q&A",
    ],
    bullet_size=19
)

# ── SLIDE 3 · Problem Statement ──────────────────────────────────────────────
bullet_slide(
    "The Drug Discovery Challenge",
    "Why computational pre-screening matters",
    [
        "💊  ~90 % of drug candidates that enter clinical trials ultimately fail.",
        "⏱  Traditional hit-to-lead optimisation takes 3–5 years and $100 M+.",
        "🔬  Computational chemists jump between 10+ isolated tools — wasting weeks.",
        "☠  Hidden toxicity (liver damage, cardiac arrest risk) discovered too late.",
        "📉  Poor oral bioavailability discovered only after in-vivo experiments.",
        "💡  Solution: One unified platform that catches problems in MINUTES.",
    ],
    bullet_size=21
)

# ── SLIDE 4 · Platform Overview ──────────────────────────────────────────────
bullet_slide(
    "ChemoFilter — Platform Overview",
    "A unified in-silico drug screening ecosystem",
    [
        "🧪  Input: Any chemical structure in SMILES notation",
        "⚙  Engine: RDKit (open-source cheminformatics) + 21 analysis modules",
        "📊  Output: 380+ physicochemical properties, drug-likeness scores, toxicity flags",
        "🤖  AI Layer: Anthropic Claude explains results in plain language",
        "📁  Batch Mode: Screen thousands of compounds from a CSV file",
        "🌐  Deployment: Streamlit web app — no installation needed",
        "🔒  Privacy: All computation is local — proprietary structures stay private",
    ],
    bullet_size=20
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 1 · ADMET & Physicochemical Properties
# ════════════════════════════════════════════════════════════════════════════
section_divider(1,
    "ADMET & Physicochemical Properties",
    "Defining the key metrics used to evaluate every molecule")

# SLIDE: ADMET explained
definition_slide(
    "What is ADMET?",
    "The five pillars of drug behaviour inside the human body",
    [
        ("A — Absorption",
         "How much of an orally administered drug reaches the bloodstream. "
         "Poor absorption = drug never works. Assessed via LogP, TPSA, solubility."),
        ("D — Distribution",
         "How the drug spreads through tissues, organs, and plasma after absorption. "
         "Governed by protein binding and lipophilicity (LogP)."),
        ("M — Metabolism",
         "How the body chemically modifies the drug, mainly in the liver via CYP450 enzymes. "
         "Rapid metabolism can inactivate a drug before it reaches its target."),
        ("E — Excretion",
         "How the drug and its metabolites are removed from the body (kidney urine, bile). "
         "Poor excretion leads to toxic accumulation."),
        ("T — Toxicity",
         "Harmful effects on cells, organs, or DNA. Includes liver toxicity (hepatotoxicity), "
         "kidney damage (nephrotoxicity), and cardiac risk (hERG inhibition)."),
    ]
)

# SLIDE: Molecular Weight
definition_slide(
    "Molecular Weight (MW)",
    "Physicochemical property — size of the molecule",
    [
        ("Definition",
         "The sum of atomic masses of all atoms in a molecule, expressed in "
         "Daltons (Da) or g/mol."),
        ("Why it matters",
         "Large molecules (MW > 500 Da) struggle to cross biological membranes "
         "and are poorly absorbed from the gut."),
        ("Lipinski Rule",
         "MW ≤ 500 Da for oral drug candidates (Rule of 5)."),
        ("Example",
         "Aspirin = 180 Da ✅ | Cyclosporin A = 1,203 Da ❌ (oral bioavailability poor)"),
        ("ChemoFilter check",
         "Calculated automatically by RDKit; flagged if > 500 Da with colour-coded warning."),
    ]
)

# SLIDE: LogP
definition_slide(
    "LogP — Partition Coefficient",
    "Measure of a molecule's fat-solubility vs water-solubility",
    [
        ("Definition",
         "Log₁₀ of the ratio of compound concentration in octanol vs water. "
         "Quantifies how 'oily' or 'water-loving' a molecule is."),
        ("Positive LogP",
         "Molecule prefers fat/lipid environments → crosses cell membranes easily "
         "but may be poorly soluble in blood plasma."),
        ("Negative LogP",
         "Molecule prefers water → high solubility but poor membrane penetration."),
        ("Optimal range",
         "LogP between −0.4 and +5.0 for most orally active drugs "
         "(Lipinski Rule: ≤ 5)."),
        ("Example",
         "Ibuprofen LogP ≈ 3.5 ✅ | Very lipophilic CNS drugs LogP > 5 → toxicity risk"),
    ]
)

# SLIDE: TPSA
definition_slide(
    "TPSA — Topological Polar Surface Area",
    "Predicts membrane permeability and oral absorption",
    [
        ("Definition",
         "The surface area (in Ų) of all polar atoms (O, N, S) in a molecule "
         "including attached hydrogens."),
        ("Why it matters",
         "Large polar surface = molecule cannot easily dissolve through the "
         "non-polar lipid bilayer of cell membranes."),
        ("Drug absorption rule",
         "TPSA ≤ 140 Ų → adequate intestinal absorption. "
         "TPSA ≤ 90 Ų → good CNS penetration (crosses blood-brain barrier)."),
        ("Calculation",
         "Computed purely from 2D topology (atom types + neighbours) "
         "— fast and reproducible with RDKit."),
        ("Example",
         "Paracetamol TPSA ≈ 49 Ų ✅ (well-absorbed) | Glucose TPSA ≈ 110 Ų "
         "(poorly absorbed, needs transporters)"),
    ]
)

# SLIDE: HBD / HBA
definition_slide(
    "HBD & HBA — Hydrogen Bond Donors & Acceptors",
    "Polarity descriptors governing solubility and permeability",
    [
        ("HBD — Donor",
         "Atoms with an N–H or O–H bond that can donate a hydrogen in a "
         "hydrogen bond. More donors → more polar → harder to cross membranes."),
        ("HBA — Acceptor",
         "Electronegative atoms (O, N, F) that can accept a hydrogen bond. "
         "Excess acceptors → high polarity → poor absorption."),
        ("Lipinski Limits",
         "HBD ≤ 5 and HBA ≤ 10 for orally bioavailable drugs (Rule of 5)."),
        ("Physical intuition",
         "Think of H-bonds as 'sticky patches'. Too many patches = drug "
         "sticks to water molecules and can't enter cells."),
        ("Example",
         "Aspirin: HBD=1, HBA=3 ✅ | Peptide drugs: HBD>10 → poor oral absorption ❌"),
    ]
)

# SLIDE: Rotatable Bonds & QED
definition_slide(
    "Rotatable Bonds & QED Score",
    "Flexibility and overall drug-likeness metrics",
    [
        ("Rotatable Bonds",
         "Single bonds (not in rings) that allow the molecule to change shape. "
         "More rotatable bonds → more flexible → harder to adopt one stable binding pose."),
        ("Rotatable Bond limit",
         "≤ 10 rotatable bonds preferred. Very flexible molecules (> 10) show "
         "poor oral bioavailability (Veber rules)."),
        ("QED — Definition",
         "Quantitative Estimate of Drug-likeness. A single score 0–1 combining "
         "MW, LogP, HBD, HBA, TPSA, rotatable bonds, aromaticity, and alerts."),
        ("QED interpretation",
         "QED > 0.67 = drug-like (top 10% of approved drugs). "
         "QED < 0.34 = poor drug-likeness. "),
        ("Example",
         "Ibuprofen QED ≈ 0.73 ✅ | Random synthetic compound QED ≈ 0.2 ❌"),
    ]
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 2 · Drug-likeness Filters
# ════════════════════════════════════════════════════════════════════════════
section_divider(2,
    "Drug-likeness Filters",
    "Rule-based and pattern-based screens to eliminate poor candidates")

definition_slide(
    "Lipinski's Rule of 5",
    "The gold-standard filter for oral bioavailability — Pfizer, 1997",
    [
        ("Origin",
         "Proposed by Christopher Lipinski at Pfizer after analysing 2,245 "
         "orally active drugs in the World Drug Index."),
        ("The 5 rules",
         "MW ≤ 500 Da | LogP ≤ 5 | HBD ≤ 5 | HBA ≤ 10 | (Rotatable bonds ≤ 10)"),
        ("Interpretation",
         "A compound violating ≥ 2 rules is very likely to have poor oral "
         "absorption or membrane permeability."),
        ("Important caveat",
         "Exceptions exist: natural product scaffolds, prodrugs, and "
         "transporter substrates often violate Ro5 yet are still oral drugs."),
        ("ChemoFilter action",
         "Computes all four Lipinski parameters and flags any violations "
         "with a coloured badge (Green=pass, Red=fail)."),
    ]
)

definition_slide(
    "Muegge Filter",
    "Extended drug-likeness criteria from Bayer — broader than Lipinski",
    [
        ("Origin",
         "Developed by Ingo Muegge at Bayer AG using trade drug databases; "
         "more liberal than Lipinski for lead optimisation."),
        ("Criteria",
         "MW 200–600 Da | LogP −2 to +5 | TPSA ≤ 150 Ų | "
         "Rotatable bonds ≤ 15 | Rings ≤ 7 | HBD ≤ 5 | HBA ≤ 10"),
        ("Advantage over Lipinski",
         "Allows slightly larger, more complex molecules suitable "
         "for modern medicinal chemistry fragment-to-lead campaigns."),
        ("Use case",
         "Applied when screening for 'lead-like' compounds — molecules "
         "with higher potency potential even at the cost of some oral absorption."),
        ("ChemoFilter action",
         "Runs Muegge checks in parallel with Lipinski — both pass/fail "
         "summaries shown in the Drug-likeness tab."),
    ]
)

definition_slide(
    "PAINS — Pan-Assay Interference Compounds",
    "Eliminating false positives from high-throughput screening",
    [
        ("What are PAINS?",
         "Chemical substructures (functional groups) that produce 'hits' "
         "in biological assays not because they bind the target, but due to "
         "non-specific interference (fluorescence, redox activity, aggregation)."),
        ("The problem",
         "A PAINS compound appears to be a drug hit in every assay tested "
         "— it creates false leads, wasting enormous resources."),
        ("How detected",
         "SMARTS pattern matching against 480 known PAINS substructures "
         "(Baell & Holloway, J. Med. Chem. 2010)."),
        ("Common PAINS",
         "Rhodanines, catechols, quinones, Michael acceptors, "
         "frequent-hitter scaffolds."),
        ("ChemoFilter action",
         "Any PAINS match immediately flags the compound as 'Unreliable for "
         "screening' — prevents wasted wet-lab follow-up."),
    ]
)

definition_slide(
    "Brenk Structural Alerts",
    "Flagging reactive and unstable chemical groups",
    [
        ("Origin",
         "Published by Brenk et al. (ChemMedChem, 2008) — "
         "105 SMARTS patterns for problematic functional groups."),
        ("What they flag",
         "Metabolically unstable groups, reactive electrophiles, "
         "genotoxic functionalities, poor pharmacokinetic groups."),
        ("Examples",
         "Acyl halides | Epoxides | Aldehydes | Nitroaromatics | "
         "Peroxides | Thiols | Diazonium salts"),
        ("Difference from PAINS",
         "PAINS = assay interference. Brenk = chemical instability "
         "or direct toxicity. Both are important to check."),
        ("ChemoFilter action",
         "Lists all matched Brenk alerts by name with the offending "
         "substructure highlighted, allowing targeted structural modification."),
    ]
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 3 · Toxicity Indicators
# ════════════════════════════════════════════════════════════════════════════
section_divider(3,
    "Toxicity Indicators",
    "Organ-specific and genetic toxicity terms defined")

definition_slide(
    "Hepatotoxicity & Nephrotoxicity",
    "Liver and kidney damage — the two most common causes of drug withdrawal",
    [
        ("Hepatotoxicity",
         "Drug-induced liver injury (DILI). The liver processes most drugs "
         "via CYP450 enzymes; reactive metabolites can destroy hepatocytes."),
        ("Hepatotoxicity markers",
         "Structural alerts: acetaminophen-like quinones, hydrazines, "
         "aromatic amines → bioactivated to toxic metabolites."),
        ("Nephrotoxicity",
         "Kidney damage caused by a drug or its metabolites. Kidneys filter "
         "blood and concentrate drugs — high local exposure damages tubules."),
        ("Nephrotoxicity markers",
         "High MW polar compounds, heavy metal complexes, cyclosporin-like "
         "structures that accumulate in renal tubular cells."),
        ("Clinical impact",
         "~30% of post-market drug withdrawals are due to DILI. "
         "ChemoFilter flags both using rule-based structure alerts."),
    ]
)

definition_slide(
    "Ames Mutagenicity Test",
    "Predicting DNA-damaging potential from molecular structure",
    [
        ("What is the Ames test?",
         "A bacterial assay (Salmonella typhimurium) used since 1973 "
         "to detect chemicals that cause mutations in DNA."),
        ("How it works (wet lab)",
         "Bacteria with a histidine mutation are exposed to the test compound; "
         "revertant colonies indicate the compound can mutate DNA."),
        ("In-silico prediction",
         "Cheminformatics models trained on thousands of experimental Ames "
         "results predict mutagenicity from SMARTS patterns."),
        ("Key mutagenic groups",
         "Nitroaromatics | Primary aromatic amines | Alkylating agents | "
         "Azo compounds | Acridines"),
        ("Regulatory importance",
         "ICH S2(R1) guidelines mandate Ames testing for all new drugs. "
         "ChemoFilter uses structural alerts as an early warning before lab work."),
    ]
)

definition_slide(
    "hERG Inhibition",
    "Predicting dangerous cardiac side effects",
    [
        ("What is hERG?",
         "Human Ether-à-go-go Related Gene — encodes a cardiac potassium "
         "ion channel (Kv11.1) essential for normal heart rhythm."),
        ("Why blocking hERG is dangerous",
         "Inhibiting hERG prolongs the QT interval on ECG → potentially fatal "
         "arrhythmia (Torsades de Pointes → ventricular fibrillation → death)."),
        ("Structural features that block hERG",
         "Basic nitrogen (pKa 8–9) combined with aromatic rings and "
         "hydrophobic bulk — many drugs have this pharmacophore accidentally."),
        ("Famous withdrawals",
         "Cisapride (stomach drug), Terfenadine (antihistamine), "
         "Astemizole — all withdrawn due to hERG-related deaths."),
        ("ChemoFilter action",
         "Computes hERG risk score from molecular features; alerts chemist "
         "to cardiac liability before any animal testing."),
    ]
)

definition_slide(
    "Blood-Brain Barrier (BBB) Permeability",
    "Predicting whether a drug can enter the brain",
    [
        ("What is the BBB?",
         "A highly selective semi-permeable membrane separating circulating "
         "blood from the brain's extracellular fluid — formed by tight junctions."),
        ("CNS drug design",
         "CNS drugs MUST cross the BBB. Peripheral drugs must NOT cross "
         "(to avoid neurological side effects)."),
        ("Molecular rules for BBB crossing",
         "MW < 450 Da | LogP 1–3 | TPSA < 90 Ų | HBD ≤ 3 | "
         "Not a P-gp efflux transporter substrate."),
        ("Example",
         "Diazepam (Valium) BBB+ ✅ — small, lipophilic | "
         "Gentamicin (antibiotic) BBB− ✅ — stays peripheral, no CNS toxicity"),
        ("ChemoFilter action",
         "Predicts BBB+ / BBB− status using TPSA, LogP, and MW thresholds; "
         "critical for both CNS drug design and safety profiling."),
    ]
)

definition_slide(
    "CYP450 Metabolism",
    "How the liver's enzyme family transforms drug molecules",
    [
        ("What are CYP450s?",
         "Cytochrome P450 enzymes — a superfamily of ~57 human enzymes "
         "responsible for oxidising ~75% of all marketed drugs."),
        ("Key isoforms",
         "CYP3A4 (handles ~50% of drugs) | CYP2D6 | CYP2C9 | "
         "CYP2C19 | CYP1A2 — each with different substrate preferences."),
        ("Inhibition problem",
         "Drug A inhibits CYP3A4 → Drug B (also CYP3A4 substrate) levels "
         "spike dangerously → drug–drug interaction (DDI)."),
        ("Induction problem",
         "Drug A induces CYP3A4 production → Drug B metabolised too fast "
         "→ therapeutic failure (e.g., rifampicin reduces contraceptive efficacy)."),
        ("ChemoFilter action",
         "Uses structural rules to predict whether a compound inhibits or "
         "induces major CYP isoforms — essential for DDI risk assessment."),
    ]
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 4 · Structural Analysis
# ════════════════════════════════════════════════════════════════════════════
section_divider(4,
    "Structural Analysis",
    "Chemical notation, representation, and substructure concepts")

definition_slide(
    "SMILES Notation",
    "The universal language for writing chemical structures as text",
    [
        ("What is SMILES?",
         "Simplified Molecular Input Line Entry System — a text-based "
         "notation encoding a molecule's atoms, bonds, and topology."),
        ("How it works",
         "Atoms represented by element symbols. Bonds: single (−), double "
         "(=), triple (#). Branches in parentheses ()."),
        ("Example",
         "Aspirin = CC(=O)Oc1ccccc1C(=O)O | "
         "Ethanol = CCO | Benzene = c1ccccc1"),
        ("Why important",
         "Enables computers to read, store, and compare chemical structures "
         "as simple strings — the input format for ChemoFilter."),
        ("ChemoFilter usage",
         "Paste any SMILES string → RDKit parses it → all 380+ properties "
         "calculated instantly. Invalid SMILES flagged with error message."),
    ]
)

definition_slide(
    "InChI Keys",
    "A standardised, fixed-length chemical identifier",
    [
        ("What is InChI?",
         "International Chemical Identifier — a IUPAC standard string "
         "representation of a molecule's structure (IUPAC, 2005)."),
        ("What is InChIKey?",
         "A 27-character fixed-length hash of the InChI string designed "
         "for database search and deduplication."),
        ("Example",
         "Aspirin InChIKey = BSYNRYMUTXBXSQ-UHFFFAOYSA-N"),
        ("Advantage over SMILES",
         "InChIKey is unique per molecule — no matter how you draw a structure, "
         "the InChIKey is always the same (canonical). SMILES can vary."),
        ("ChemoFilter usage",
         "Converts SMILES to InChIKey for cross-referencing against DrugBank, "
         "PubChem, and ChEMBL databases to find known analogues."),
    ]
)

definition_slide(
    "Scaffold Analysis",
    "Identifying the core ring system shared by a drug class",
    [
        ("What is a scaffold?",
         "The central ring system (core structure) of a drug molecule "
         "after removing all side chains — defines the pharmacophoric 'skeleton'."),
        ("Bemis–Murcko scaffold",
         "The most common scaffold definition: retains all ring systems "
         "and the linker atoms connecting them. Named after Bemis & Murcko, 1996."),
        ("Why scaffolds matter",
         "Drugs with the same scaffold often share mechanism of action, "
         "toxicity profiles, and patent status — vital for IP and safety."),
        ("Scaffold diversity",
         "A screening library with high scaffold diversity → better chance "
         "of finding a hit for a novel target."),
        ("ChemoFilter usage",
         "Extracts Bemis–Murcko scaffold via RDKit; plots scaffold frequency "
         "histograms for batch library analysis — shows structural clustering."),
    ]
)

definition_slide(
    "Structural Alerts & Toxicophores",
    "Molecular substructures predictive of toxicity",
    [
        ("Structural Alert",
         "A specific chemical substructure (defined by a SMARTS pattern) "
         "known to be associated with toxicity, reactivity, or poor drug behaviour."),
        ("Toxicophore",
         "A pharmacophore for toxicity — a spatial arrangement of features "
         "that causes a toxic interaction with a biological macromolecule."),
        ("Examples",
         "Nitroaromatic group → mutagenicity | "
         "Quinone → oxidative stress | "
         "Epoxide → DNA alkylation"),
        ("SMARTS patterns",
         "SMARTS = SMILES ARbitrary Target Specification — a pattern language "
         "for substructure searching (superset of SMILES, supports wildcards)."),
        ("ChemoFilter usage",
         "Checks every molecule against 480 PAINS + 105 Brenk SMARTS patterns "
         "+ custom toxicophore library — all flagged with names and explanations."),
    ]
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 5 · Core Capabilities
# ════════════════════════════════════════════════════════════════════════════
section_divider(5,
    "Core Platform Capabilities",
    "What ChemoFilter can do for your research")

two_column_slide(
    "Analysis Capabilities",
    "380+ properties computed per molecule",
    left_header="Property Calculations",
    right_header="Screening & Filters",
    left_items=[
        "✅  Molecular Weight, Formula, SMILES",
        "✅  LogP, LogD, LogS (solubility)",
        "✅  TPSA, PSA, Fsp3 (sp3 fraction)",
        "✅  HBD, HBA, Rotatable Bonds",
        "✅  Ring count, Aromatic rings",
        "✅  QED Score (0–1 drug-likeness)",
        "✅  Molar Refractivity (size + polarity)",
        "✅  Chiral centres, stereocentres",
    ],
    right_items=[
        "🔍  Lipinski Rule of 5",
        "🔍  Muegge lead-likeness filter",
        "🔍  Veber oral bioavailability rules",
        "🔍  PAINS detection (480 patterns)",
        "🔍  Brenk alert detection (105 patterns)",
        "🔍  hERG inhibition prediction",
        "🔍  BBB permeability prediction",
        "🔍  CYP450 inhibition flags",
    ]
)

two_column_slide(
    "Toxicity & AI Capabilities",
    "Safety profiling and intelligent explanations",
    left_header="Toxicity Prediction",
    right_header="AI & Visualisation",
    left_items=[
        "☢  Ames mutagenicity prediction",
        "☢  Hepatotoxicity structural alert",
        "☢  Nephrotoxicity flagging",
        "☢  Reactive metabolite detection",
        "☢  Genotoxicity (DNA-reactive groups)",
        "☢  Skin sensitisation alerts",
        "☢  Environmental toxicity flags",
    ],
    right_items=[
        "🤖  Claude AI plain-language explanation",
        "🤖  Drug analogue suggestions",
        "🤖  Repurposing opportunities",
        "📊  Radar chart (property profile)",
        "📊  3D conformer viewer",
        "📊  Scaffold clustering plot",
        "📊  Batch CSV report export",
    ]
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 6 · ChemoScore
# ════════════════════════════════════════════════════════════════════════════
section_divider(6,
    "ChemoScore — The Unified Drug Quality Score",
    "A single number summarising all analysis results")

bullet_slide(
    "ChemoScore — Methodology",
    "How the overall drug quality grade is calculated",
    [
        "📐  ChemoScore = Weighted sum of four sub-scores:",
        "     ① Structural Score    — scaffold quality, alerts, PAINS",
        "     ② PhysChem Score     — MW, LogP, TPSA within optimal ranges",
        "     ③ Safety Score        — hERG, Ames, BBB, CYP flags",
        "     ④ Synthesis Score     — synthetic accessibility, complexity",
        "⚖  Researchers adjust weights via interactive sliders (e.g., safety × 2 for CNS drugs)",
        "🏅  Grades: A (≥ 80) · B (60–79) · C (40–59) · F (< 40)",
        "🚀  Allows instant ranking of 1,000+ compounds — top hits identified in seconds",
    ],
    bullet_size=20
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 7 · Computational Workflow
# ════════════════════════════════════════════════════════════════════════════
section_divider(7,
    "Computational Workflow",
    "Step-by-step: how a molecule is analysed")

bullet_slide(
    "How ChemoFilter Analyses a Molecule",
    "From SMILES string to actionable report",
    [
        "STEP 1 · Input   — Paste SMILES string (single) or upload CSV (batch)",
        "STEP 2 · Parse   — RDKit converts SMILES → RDKit Mol object (validates structure)",
        "STEP 3 · Feature Engine  — Calculates MW, LogP, TPSA, HBD, HBA, QED, Fsp3 …",
        "STEP 4 · Filter Engine   — Checks Lipinski / Muegge / PAINS / Brenk rules",
        "STEP 5 · Toxicity Engine — Evaluates hepatotoxicity, hERG, Ames, CYP patterns",
        "STEP 6 · ChemoScore      — Combines all scores with user-defined weights → Grade",
        "STEP 7 · AI Explanation  — Claude API generates chemist-friendly summary",
        "STEP 8 · Export          — Download full report as CSV / JSON / text",
    ],
    bullet_size=19
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 8 · AI Integration
# ════════════════════════════════════════════════════════════════════════════
section_divider(8,
    "AI Integration — Claude by Anthropic",
    "Making results understandable for every chemist")

bullet_slide(
    "AI-Powered Explanations",
    "How Anthropic Claude enhances ChemoFilter",
    [
        "🤖  Natural Language Summaries: 'This compound has good oral absorption but cardiac risk due to its basic nitrogen group adjacent to the aromatic ring...'",
        "💊  Drug Analogues: Identifies structurally similar approved drugs in market → accelerates lead selection",
        "🔄  Repurposing Suggestions: Finds existing drugs that might work for new disease targets",
        "⚠  Risk Narratives: Translates toxicity flags into plain English for interdisciplinary teams",
        "🛡  Graceful Fallback: If API unavailable, ChemoFilter continues working — AI layer is optional enhancement",
        "🔒  Privacy: Only non-proprietary structural fragments sent to API — full SMILES never transmitted",
    ],
    bullet_size=19
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 9 · Unique Features
# ════════════════════════════════════════════════════════════════════════════
section_divider(9,
    "What Makes ChemoFilter Unique",
    "Comparing to traditional workflows and existing tools")

two_column_slide(
    "ChemoFilter vs Traditional Workflow",
    "Why the unified approach wins",
    left_header="Traditional Workflow",
    right_header="ChemoFilter",
    left_items=[
        "❌  10+ separate tools (PAINS server,\n    LogP calc, TPSA tool, 3D viewer …)",
        "❌  Manual copy-paste between tools",
        "❌  Days to screen 100 compounds",
        "❌  Results in different file formats",
        "❌  No unified scoring system",
        "❌  No AI explanation layer",
        "❌  Requires expert chemist to interpret",
    ],
    right_items=[
        "✅  Single unified web application",
        "✅  Automated pipeline, zero manual steps",
        "✅  Minutes to screen 1,000+ compounds",
        "✅  Unified CSV / JSON export",
        "✅  ChemoScore grade (A/B/C/F)",
        "✅  Claude AI plain-language summary",
        "✅  Accessible to entire research team",
    ]
)

# ════════════════════════════════════════════════════════════════════════════
# SECTION 10 · Future Scope
# ════════════════════════════════════════════════════════════════════════════
section_divider(10,
    "Future Scope & Roadmap",
    "Planned enhancements to extend ChemoFilter's capabilities")

bullet_slide(
    "Future Scope — Phase 1 (Near Term)",
    "Extensions within the current architecture",
    [
        "🧬  Machine Learning ADMET Models — trained on ChEMBL/BindingDB bioactivity data",
        "🏗  3D Protein Docking Simulation — predict binding pose at target active site",
        "🔗  DrugBank & PubChem API Integration — cross-reference 100M+ real compounds",
        "📐  Retrosynthesis Analysis — suggest optimal synthetic route to the target molecule",
        "📊  Interactive Property Heatmaps — visual comparison across entire screening library",
        "⚗  Experimental Data Upload — connect physical lab assay results back into platform",
    ],
    bullet_size=19
)

bullet_slide(
    "Future Scope — Phase 2 (Long Term)",
    "Advanced research and regulatory capabilities",
    [
        "🧠  Explainable AI (XAI) — show WHY each prediction was made, not just what",
        "🎯  Multi-Target Screening — efficacy + safety across multiple disease targets simultaneously",
        "👥  Team Collaboration — shared lab notebooks, annotation, project management",
        "🏛  FDA Regulatory Compliance Checker — automatic ICH guideline assessment",
        "⚖  Patent Similarity Analysis — avoid infringing existing IP before synthesis",
        "🌍  Personalised Toxicity Profiles — account for ethnic/genetic CYP450 polymorphisms",
        "🔬  IoT Lab Connectivity — link to HPLC, mass spectrometry instruments in real time",
    ],
    bullet_size=18
)

# ── References ───────────────────────────────────────────────────────────────
section_divider("",
    "Scientific References",
    "Peer-reviewed literature underpinning ChemoFilter's algorithms")

bullet_slide(
    "Key References",
    "Scientific basis of ChemoFilter's rules and predictions",
    [
        "Lipinski CA et al. (1997). Experimental and computational approaches to estimate "
        "solubility and permeability. Adv. Drug Deliv. Rev. 23, 3–25.",
        "Baell JB & Holloway GA (2010). New substructure filters for removal of PAINS "
        "from screening libraries. J. Med. Chem. 53, 2719–2740.",
        "Brenk R et al. (2008). Lessons learned from assembling screening libraries "
        "for drug discovery. ChemMedChem 3, 435–444.",
        "Muegge I (2003). Selection criteria for drug-like compounds. Med. Res. Rev. 23, 302–321.",
        "Ertl P & Schuffenhauer A (2009). Estimation of synthetic accessibility score "
        "of drug-like molecules. J. Cheminformatics 1, 8.",
        "Bemis GW & Murcko MA (1996). The properties of known drugs. "
        "J. Med. Chem. 39, 2887–2893.",
        "RDKit: Open-source cheminformatics. https://www.rdkit.org",
    ],
    bullet_size=15,
    bullet_color=LIGHT_GRAY
)

# ── Thank You ─────────────────────────────────────────────────────────────────
thank_you_slide()

# ── Save ──────────────────────────────────────────────────────────────────────
OUTPUT = "ChemoFilter_Detailed_Presentation.pptx"
prs.save(OUTPUT)
print(f"✅  Presentation saved: {OUTPUT}  ({len(prs.slides)} slides)")
